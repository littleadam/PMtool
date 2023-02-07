import gspread
from oauth2client.service_account import ServiceAccountCredentials
import pptx

# Define the credentials for the Google Sheets API
scope = ['https://spreadsheets.google.com/feeds',
         'https://www.googleapis.com/auth/drive']
creds = ServiceAccountCredentials.from_json_keyfile_name('service_account.json', scope)
client = gspread.authorize(creds)

# Open the Google Sheet and get the data
sheet = client.open('Your_Sheet_Name').sheet1
data = sheet.get_all_values()

# Create a new PowerPoint presentation
presentation = pptx.Presentation()

# Add a title slide
title_slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = 'Google Sheets Data'

# Add a slide for each row of data
for i, row in enumerate(data):
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    shapes = slide.shapes

    # Add the data from each column as text boxes on the slide
    for j, item in enumerate(row):
        left = (j * pptx.util.Inches(2))
        top = (i * pptx.util.Inches(0.5))
        width = pptx.util.Inches(2)
        height = pptx.util.Inches(0.5)
        shape = shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
        text_frame.text = item

# Save the PowerPoint presentation
presentation.save('google_sheets_data.pptx')
